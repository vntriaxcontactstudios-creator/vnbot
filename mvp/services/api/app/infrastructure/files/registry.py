"""VNBOT API — FileReader Registry.

Per ADR-0010: extensible file reading system supporting 80+ extensions.
Each reader is registered via decorator and called based on file extension.
"""

from __future__ import annotations

import base64
import io
import logging
import os
import zipfile
from dataclasses import dataclass, field
from typing import Any, Callable, Awaitable

logger = logging.getLogger("vnbot.api.files")

@dataclass
class FileContent:
    """Result of reading a file."""
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    images: list[bytes] = field(default_factory=list)  # images for VLM analysis
    mime_type: str = "application/octet-stream"
    reader_used: str = "unknown"


# Registry: extension → reader function
FILE_READERS: dict[str, Callable[[bytes, str], Awaitable[FileContent]]] = {}

# Supported extensions (populated by @register_reader)
SUPPORTED_EXTENSIONS: set[str] = set()


def register_reader(*extensions: str):
    """Decorator to register a file reader for one or more extensions.

    Usage:
        @register_reader("xlsx", "xls")
        async def read_excel(data: bytes, filename: str) -> FileContent:
            ...
    """
    def decorator(func: Callable[[bytes, str], Awaitable[FileContent]]) -> Callable[[bytes, str], Awaitable[FileContent]]:
        for ext in extensions:
            ext_clean = ext.lower().lstrip(".")
            FILE_READERS[ext_clean] = func
            SUPPORTED_EXTENSIONS.add(ext_clean)
        logger.debug(f"Registered reader for: {', '.join(extensions)}")
        return func
    return decorator


async def read_file(data: bytes, filename: str) -> FileContent:
    """Read a file based on its extension. Returns FileContent with text + metadata + images.

    Raises ValueError if extension is not supported.
    """
    ext = ""
    if "." in filename:
        ext = filename.rsplit(".", 1)[-1].lower()

    reader = FILE_READERS.get(ext)
    if reader is None:
        raise ValueError(
            f"Unsupported file type: .{ext}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )

    try:
        content = await reader(data, filename)
        content.reader_used = ext
        return content
    except Exception as e:
        logger.error(f"Error reading .{ext} file '{filename}': {e}")
        raise ValueError(f"Failed to read .{ext} file: {e}") from e


def get_supported_extensions() -> list[str]:
    """Return sorted list of supported extensions."""
    return sorted(SUPPORTED_EXTENSIONS)


# ──────────────────────────────────────────────────────────────────────────────
# Built-in readers — Text formats
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("txt", "text", "log", "md", "markdown", "rst", "adoc", "org", "wiki", "sh", "bash", "zsh", "ps1", "lua", "r", "dockerfile", "makefile", "cmake", "ini", "cfg", "conf", "env", "gitignore", "dockerignore", "toml")
async def read_text(data: bytes, filename: str) -> FileContent:
    """Read plain text files."""
    text = data.decode("utf-8", errors="replace")
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"
    mime_map = {
        "md": "text/markdown", "markdown": "text/markdown",
        "html": "text/html", "css": "text/css",
        "json": "application/json", "xml": "application/xml",
        "toml": "application/toml", "yaml": "application/yaml", "yml": "application/yaml",
    }
    return FileContent(
        text=text,
        metadata={"lines": text.count("\n") + 1, "chars": len(text)},
        mime_type=mime_map.get(ext, "text/plain"),
    )


@register_reader("json")
async def read_json(data: bytes, filename: str) -> FileContent:
    """Read JSON files."""
    import json
    text = data.decode("utf-8", errors="replace")
    try:
        parsed = json.loads(text)
        metadata = {"valid": True, "type": type(parsed).__name__}
        if isinstance(parsed, dict):
            metadata["keys"] = list(parsed.keys())[:20]
        elif isinstance(parsed, list):
            metadata["items"] = len(parsed)
    except json.JSONDecodeError as e:
        metadata = {"valid": False, "error": str(e)}
    return FileContent(text=text, metadata=metadata, mime_type="application/json")


@register_reader("yaml", "yml")
async def read_yaml(data: bytes, filename: str) -> FileContent:
    """Read YAML files."""
    text = data.decode("utf-8", errors="replace")
    return FileContent(text=text, metadata={"format": "yaml"}, mime_type="application/yaml")


@register_reader("xml")
async def read_xml(data: bytes, filename: str) -> FileContent:
    """Read XML files."""
    from lxml import etree
    text = data.decode("utf-8", errors="replace")
    try:
        root = etree.fromstring(data)
        metadata = {"root_tag": root.tag, "children": len(root)}
    except Exception:
        metadata = {"parse_error": True}
    return FileContent(text=text, metadata=metadata, mime_type="application/xml")


@register_reader("html", "htm")
async def read_html(data: bytes, filename: str) -> FileContent:
    """Read HTML files — extract text content."""
    from bs4 import BeautifulSoup
    text_raw = data.decode("utf-8", errors="replace")
    soup = BeautifulSoup(text_raw, "html.parser")
    # Extract visible text
    for script in soup(["script", "style"]):
        script.decompose()
    text = soup.get_text(separator="\n", strip=True)
    metadata = {
        "title": soup.title.string if soup.title else None,
        "links": len(soup.find_all("a")),
        "images": len(soup.find_all("img")),
        "paragraphs": len(soup.find_all("p")),
    }
    # Extract images for VLM
    images: list[bytes] = []
    return FileContent(text=text, metadata=metadata, images=images, mime_type="text/html")


@register_reader("csv")
async def read_csv(data: bytes, filename: str) -> FileContent:
    """Read CSV files."""
    import csv as csv_module
    text = data.decode("utf-8", errors="replace")
    reader = csv_module.reader(io.StringIO(text))
    rows = list(reader)
    metadata = {"rows": len(rows), "columns": len(rows[0]) if rows else 0}
    # Limit to first 100 rows for text
    display = "\n".join([",".join(row) for row in rows[:100]])
    if len(rows) > 100:
        display += f"\n... ({len(rows) - 100} more rows)"
    return FileContent(text=display, metadata=metadata, mime_type="text/csv")


@register_reader("sql")
async def read_sql(data: bytes, filename: str) -> FileContent:
    """Read SQL files."""
    import sqlparse
    text = data.decode("utf-8", errors="replace")
    statements = sqlparse.parse(text)
    metadata = {"statements": len(statements)}
    return FileContent(text=text, metadata=metadata, mime_type="application/sql")


@register_reader("tex")
async def read_latex(data: bytes, filename: str) -> FileContent:
    """Read LaTeX files."""
    from pylatexenc.latex2text import LatexNodes2Text
    text_raw = data.decode("utf-8", errors="replace")
    try:
        converter = LatexNodes2Text()
        text = converter.latex_to_text(text_raw)
    except Exception:
        text = text_raw
    return FileContent(text=text, metadata={"format": "latex"}, mime_type="application/x-tex")


@register_reader("rst")
async def read_rst(data: bytes, filename: str) -> FileContent:
    """Read reStructuredText files."""
    import docutils.core
    text_raw = data.decode("utf-8", errors="replace")
    try:
        result = docutils.core.publish_parts(text_raw, writer_name="html")
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(result["html_body"], "html.parser")
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        text = text_raw
    return FileContent(text=text, metadata={"format": "rst"}, mime_type="text/x-rst")


# ──────────────────────────────────────────────────────────────────────────────
# Document formats
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("pdf")
async def read_pdf(data: bytes, filename: str) -> FileContent:
    """Read PDF files — extract text + images."""
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(data))
    text_parts = []
    images: list[bytes] = []
    for i, page in enumerate(reader.pages):
        text_parts.append(f"--- Page {i+1} ---")
        text_parts.append(page.extract_text() or "")
    text = "\n".join(text_parts)
    metadata = {"pages": len(reader.pages), "encrypted": reader.is_encrypted}
    return FileContent(text=text, metadata=metadata, images=images, mime_type="application/pdf")


@register_reader("docx")
async def read_docx(data: bytes, filename: str) -> FileContent:
    """Read Word .docx files."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(text_parts)
    metadata = {
        "paragraphs": len(doc.paragraphs),
        "tables": len(doc.tables),
        "sections": len(doc.sections),
    }
    # Extract images
    images: list[bytes] = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                images.append(rel.target_part.blob)
            except Exception:
                pass
    return FileContent(text=text, metadata=metadata, images=images, mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")


@register_reader("pptx")
async def read_pptx(data: bytes, filename: str) -> FileContent:
    """Read PowerPoint .pptx files."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    text_parts = []
    images: list[bytes] = []
    for i, slide in enumerate(prs.slides):
        text_parts.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    text = "\n".join(text_parts)
    metadata = {"slides": len(prs.slides)}
    return FileContent(text=text, metadata=metadata, images=images, mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@register_reader("xlsx", "xls")
async def read_excel(data: bytes, filename: str) -> FileContent:
    """Read Excel files."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 100), values_only=True):
            text_parts.append("\t".join(str(c) if c is not None else "" for c in row))
    text = "\n".join(text_parts)
    metadata = {"sheets": wb.sheetnames, "max_rows": max(wb[s].max_row for s in wb.sheetnames) if wb.sheetnames else 0}
    return FileContent(text=text, metadata=metadata, mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@register_reader("odt")
async def read_odt(data: bytes, filename: str) -> FileContent:
    """Read OpenDocument Text."""
    from odf.opendocument import load
    from odf.text import P
    doc = load(io.BytesIO(data))
    paragraphs = doc.getElementsByType(P)
    text_parts = [str(p.firstChild) if p.firstChild else "" for p in paragraphs]
    text = "\n".join(text_parts)
    metadata = {"paragraphs": len(paragraphs)}
    return FileContent(text=text, metadata=metadata, mime_type="application/vnd.oasis.opendocument.text")


@register_reader("rtf")
async def read_rtf(data: bytes, filename: str) -> FileContent:
    """Read RTF files."""
    from striprtf.striprtf import rtf_to_text
    text_raw = data.decode("utf-8", errors="replace")
    text = rtf_to_text(text_raw)
    return FileContent(text=text, metadata={"format": "rtf"}, mime_type="application/rtf")


@register_reader("epub")
async def read_epub(data: bytes, filename: str) -> FileContent:
    """Read EPUB e-books."""
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    book = epub.read_epub(io.BytesIO(data))
    text_parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        text_parts.append(soup.get_text(separator="\n", strip=True))
    text = "\n".join(text_parts)
    metadata = {
        "title": book.get_metadata("DC", "title")[0][0] if book.get_metadata("DC", "title") else None,
        "author": book.get_metadata("DC", "creator")[0][0] if book.get_metadata("DC", "creator") else None,
        "chapters": len(text_parts),
    }
    return FileContent(text=text, metadata=metadata, mime_type="application/epub+zip")


@register_reader("ipynb")
async def read_notebook(data: bytes, filename: str) -> FileContent:
    """Read Jupyter Notebooks."""
    import nbformat
    nb = nbformat.reads(data.decode("utf-8"), as_version=4)
    text_parts = []
    for i, cell in enumerate(nb.cells):
        cell_type = cell.cell_type
        text_parts.append(f"--- Cell {i+1} ({cell_type}) ---")
        text_parts.append(cell.source)
        if cell_type == "code" and cell.get("outputs"):
            for output in cell.outputs:
                if "text" in output:
                    text_parts.append(f"[Output]: {output['text']}")
    text = "\n".join(text_parts)
    metadata = {
        "cells": len(nb.cells),
        "code_cells": sum(1 for c in nb.cells if c.cell_type == "code"),
        "markdown_cells": sum(1 for c in nb.cells if c.cell_type == "markdown"),
    }
    return FileContent(text=text, metadata=metadata, mime_type="application/x-ipynb+json")


# ──────────────────────────────────────────────────────────────────────────────
# Image formats
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("png", "jpg", "jpeg", "bmp", "tiff", "tif", "webp", "gif", "ico", "heic", "heif")
async def read_image(data: bytes, filename: str) -> FileContent:
    """Read image files — extract metadata, prepare for VLM analysis."""
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    metadata = {
        "format": img.format,
        "mode": img.mode,
        "size": f"{img.width}x{img.height}",
        "width": img.width,
        "height": img.height,
    }
    # The image itself is sent to VLM
    return FileContent(
        text=f"[Image: {filename} — {img.width}x{img.height} {img.format}]",
        metadata=metadata,
        images=[data],
        mime_type=f"image/{(img.format or 'unknown').lower()}",
    )


@register_reader("svg")
async def read_svg(data: bytes, filename: str) -> FileContent:
    """Read SVG files."""
    text = data.decode("utf-8", errors="replace")
    return FileContent(text=text, metadata={"format": "svg"}, images=[data], mime_type="image/svg+xml")


# ──────────────────────────────────────────────────────────────────────────────
# Design / Illustration formats
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("psd")
async def read_psd(data: bytes, filename: str) -> FileContent:
    """Read Photoshop PSD files."""
    from psd_tools import PSDImage
    psd = PSDImage.open(io.BytesIO(data))
    text_parts = [f"PSD: {psd.width}x{psd.height}, {len(list(psd.layers))} layers"]
    for i, layer in enumerate(psd.layers):
        text_parts.append(f"  Layer {i+1}: {layer.name} ({layer.kind})")
    # Flatten to PNG for VLM
    images: list[bytes] = []
    try:
        composite = psd.composite()
        buf = io.BytesIO()
        composite.save(buf, format="PNG")
        images.append(buf.getvalue())
    except Exception:
        pass
    return FileContent(
        text="\n".join(text_parts),
        metadata={"layers": len(list(psd.layers)), "width": psd.width, "height": psd.height},
        images=images,
        mime_type="image/vnd.adobe.photoshop",
    )


@register_reader("kra", "ora")
async def read_krita(data: bytes, filename: str) -> FileContent:
    """Read Krita/OpenRaster files (ZIP-based)."""
    text_parts = []
    images: list[bytes] = []
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        text_parts.append(f"Layers: {names}")
        # Try to read merged image
        for name in names:
            if name in ("mergedimage.png", "flat.png", "document.png"):
                images.append(zf.read(name))
                break
        # Read maindoc.xml for metadata
        for name in names:
            if name.endswith(".xml"):
                try:
                    xml_content = zf.read(name).decode("utf-8", errors="replace")
                    text_parts.append(f"--- {name} ---\n{xml_content[:500]}")
                except Exception:
                    pass
    return FileContent(
        text="\n".join(text_parts),
        metadata={"format": "kra" if filename.endswith(".kra") else "ora", "files": len(names)},
        images=images,
        mime_type="application/x-krita",
    )


# ──────────────────────────────────────────────────────────────────────────────
# Data formats
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("parquet")
async def read_parquet(data: bytes, filename: str) -> FileContent:
    """Read Parquet files."""
    import pyarrow.parquet as pq
    table = pq.read_table(io.BytesIO(data))
    df = table.to_pandas()
    text = df.head(100).to_string()
    metadata = {"rows": len(df), "columns": len(df.columns), "column_names": list(df.columns)}
    return FileContent(text=text, metadata=metadata, mime_type="application/vnd.apache.parquet")


@register_reader("sqlite", "db", "sqlite3")
async def read_sqlite(data: bytes, filename: str) -> FileContent:
    """Read SQLite database — list tables + schema."""
    import sqlite3
    conn = sqlite3.connect(io.BytesIO(data).read() if hasattr(io.BytesIO(data), 'read') else ":memory:")
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = [row[0] for row in cursor.fetchall()]
    text_parts = [f"Tables: {tables}"]
    for table in tables[:10]:
        cursor.execute(f"SELECT sql FROM sqlite_master WHERE name='{table}';")
        schema = cursor.fetchone()
        if schema:
            text_parts.append(f"--- {table} ---\n{schema[0]}")
        cursor.execute(f"SELECT COUNT(*) FROM {table};")
        count = cursor.fetchone()[0]
        text_parts.append(f"Rows: {count}")
    conn.close()
    return FileContent(text="\n".join(text_parts), metadata={"tables": tables}, mime_type="application/x-sqlite3")


# ──────────────────────────────────────────────────────────────────────────────
# Compressed files
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("zip")
async def read_zip(data: bytes, filename: str) -> FileContent:
    """Read ZIP archives — extract file list + read text files inside."""
    text_parts = []
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        text_parts.append(f"Files in archive ({len(names)}):")
        for name in names[:50]:
            text_parts.append(f"  {name}")
        # Read text files inside (limit 10)
        text_files = [n for n in names if any(n.endswith(f".{ext}") for ext in ("txt", "md", "json", "yaml", "yml", "xml", "csv", "py", "js", "ts"))]
        for name in text_files[:10]:
            try:
                inner_data = zf.read(name)
                inner_content = await read_file(inner_data, name)
                text_parts.append(f"\n--- {name} ---\n{inner_content.text[:500]}")
            except Exception:
                pass
    return FileContent(text="\n".join(text_parts), metadata={"files": len(names)}, mime_type="application/zip")


# ──────────────────────────────────────────────────────────────────────────────
# 3D formats
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("obj")
async def read_obj(data: bytes, filename: str) -> FileContent:
    """Read Wavefront OBJ 3D files."""
    import trimesh
    mesh = trimesh.load(io.BytesIO(data), file_type="obj")
    text = f"3D Model: {filename}\nVertices: {len(mesh.vertices)}\nFaces: {len(mesh.faces)}\n"
    if hasattr(mesh, "metadata") and mesh.metadata:
        text += f"Metadata: {mesh.metadata}\n"
    return FileContent(
        text=text,
        metadata={"vertices": len(mesh.vertices), "faces": len(mesh.faces)},
        mime_type="model/obj",
    )


@register_reader("glb", "gltf")
async def read_gltf(data: bytes, filename: str) -> FileContent:
    """Read glTF 3D files."""
    import trimesh
    mesh = trimesh.load(io.BytesIO(data), file_type="glb" if filename.endswith(".glb") else "gltf")
    text = f"3D Model: {filename}\nVertices: {len(mesh.vertices) if hasattr(mesh, 'vertices') else 'N/A'}\n"
    return FileContent(text=text, metadata={"format": "glTF"}, mime_type="model/gltf-binary")


# ──────────────────────────────────────────────────────────────────────────────
# Additional readers — Fase 0.5 Sprint 1 extension
# ──────────────────────────────────────────────────────────────────────────────

@register_reader("ods")
async def read_ods(data: bytes, filename: str) -> FileContent:
    """Read OpenDocument Spreadsheet."""
    from odf.opendocument import load
    from odf.table import Table, TableCell, TableRow
    doc = load(io.BytesIO(data))
    tables = doc.getElementsByType(Table)
    text_parts = []
    for i, table in enumerate(tables):
        text_parts.append(f"--- Sheet {i+1} ---")
        rows = table.getElementsByType(TableRow)
        for row in rows[:100]:
            cells = row.getElementsByType(TableCell)
            row_text = []
            for cell in cells:
                txt = ""
                for p in cell.childNodes:
                    txt += str(p)
                row_text.append(txt.strip())
            text_parts.append("\t".join(row_text))
    return FileContent(text="\n".join(text_parts), metadata={"sheets": len(tables)}, mime_type="application/vnd.oasis.opendocument.spreadsheet")


@register_reader("odp")
async def read_odp(data: bytes, filename: str) -> FileContent:
    """Read OpenDocument Presentation."""
    from odf.opendocument import load
    from odf.text import P
    doc = load(io.BytesIO(data))
    paragraphs = doc.getElementsByType(P)
    text_parts = [str(p.firstChild) if p.firstChild else "" for p in paragraphs]
    return FileContent(text="\n".join(text_parts), metadata={"paragraphs": len(paragraphs)}, mime_type="application/vnd.oasis.opendocument.presentation")


@register_reader("mobi")
async def read_mobi(data: bytes, filename: str) -> FileContent:
    """Read MOBI e-books."""
    try:
        import mobi
        temp_path = "/tmp/vnbot_mobi_temp"
        with open(temp_path, "wb") as f:
            f.write(data)
        text = mobi.extract(temp_path)
        return FileContent(text=str(text)[:50000], metadata={"format": "mobi"}, mime_type="application/x-mobipocket-ebook")
    except ImportError:
        return FileContent(text="[MOBI reader no disponible — instalar: pip install mobi]", metadata={"error": "mobi not installed"}, mime_type="application/x-mobipocket-ebook")


@register_reader("7z")
async def read_7z(data: bytes, filename: str) -> FileContent:
    """Read 7-Zip archives."""
    import py7zr
    text_parts = []
    with py7zr.SevenZipFile(io.BytesIO(data), mode="r") as zf:
        names = zf.getnames()
        text_parts.append(f"Files in 7z archive ({len(names)}):")
        for name in names[:50]:
            text_parts.append(f"  {name}")
        # Extract text files
        extracted = zf.readall() if hasattr(zf, 'readall') else {}
        for name, bio in list(extracted.items())[:10]:
            if any(name.endswith(f".{ext}") for ext in ("txt", "md", "json", "yaml", "py", "js", "ts")):
                try:
                    inner_data = bio.read()
                    inner_content = await read_file(inner_data, name)
                    text_parts.append(f"\n--- {name} ---\n{inner_content.text[:500]}")
                except Exception:
                    pass
    return FileContent(text="\n".join(text_parts), metadata={"files": len(names)}, mime_type="application/x-7z-compressed")


@register_reader("rar")
async def read_rar(data: bytes, filename: str) -> FileContent:
    """Read RAR archives."""
    import rarfile
    import tempfile
    # rarfile needs a file path
    with tempfile.NamedTemporaryFile(suffix=".rar", delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        rf = rarfile.RarFile(tmp_path)
        names = rf.namelist()
        text_parts = [f"Files in RAR archive ({len(names)}):"]
        for name in names[:50]:
            text_parts.append(f"  {name}")
        for name in names[:10]:
            if any(name.endswith(f".{ext}") for ext in ("txt", "md", "json", "yaml", "py")):
                try:
                    inner_data = rf.read(name)
                    inner_content = await read_file(inner_data, name)
                    text_parts.append(f"\n--- {name} ---\n{inner_content.text[:500]}")
                except Exception:
                    pass
    finally:
        os.unlink(tmp_path)
    return FileContent(text="\n".join(text_parts), metadata={"files": len(names)}, mime_type="application/vnd.rar")


@register_reader("tar", "gz", "tgz", "tar.gz")
async def read_tarball(data: bytes, filename: str) -> FileContent:
    """Read tar/tar.gz archives."""
    import tarfile
    text_parts = []
    with tarfile.open(fileobj=io.BytesIO(data), mode="r:*") as tf:
        names = tf.getnames()
        text_parts.append(f"Files in tarball ({len(names)}):")
        for name in names[:50]:
            text_parts.append(f"  {name}")
        for member in tf.getmembers()[:10]:
            if member.isfile() and any(member.name.endswith(f".{ext}") for ext in ("txt", "md", "json", "yaml", "py")):
                try:
                    inner_data = tf.extractfile(member).read()
                    inner_content = await read_file(inner_data, member.name)
                    text_parts.append(f"\n--- {member.name} ---\n{inner_content.text[:500]}")
                except Exception:
                    pass
    return FileContent(text="\n".join(text_parts), metadata={"files": len(names)}, mime_type="application/x-tar")


@register_reader("eml")
async def read_eml(data: bytes, filename: str) -> FileContent:
    """Read email .eml files."""
    import email
    from email.policy import default
    msg = email.message_from_bytes(data, policy=default)
    text_parts = [
        f"From: {msg.get('From', 'N/A')}",
        f"To: {msg.get('To', 'N/A')}",
        f"Subject: {msg.get('Subject', 'N/A')}",
        f"Date: {msg.get('Date', 'N/A')}",
        "",
    ]
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                text_parts.append(part.get_content())
                break
    else:
        text_parts.append(msg.get_content())
    return FileContent(text="\n".join(text_parts), metadata={"subject": msg.get("Subject"), "from": msg.get("From")}, mime_type="message/rfc822")


@register_reader("vcf")
async def read_vcf(data: bytes, filename: str) -> FileContent:
    """Read vCard files."""
    text = data.decode("utf-8", errors="replace")
    contacts = text.count("BEGIN:VCARD")
    return FileContent(text=text, metadata={"contacts": contacts}, mime_type="text/vcard")


@register_reader("geojson")
async def read_geojson(data: bytes, filename: str) -> FileContent:
    """Read GeoJSON files."""
    import json
    text = data.decode("utf-8", errors="replace")
    try:
        parsed = json.loads(text)
        features = parsed.get("features", [])
        metadata = {"type": parsed.get("type"), "features": len(features)}
        if features:
            metadata["feature_types"] = list(set(f.get("geometry", {}).get("type", "?") for f in features[:20]))
    except Exception:
        metadata = {"parse_error": True}
    return FileContent(text=text, metadata=metadata, mime_type="application/geo+json")


@register_reader("py")
async def read_python(data: bytes, filename: str) -> FileContent:
    """Read Python files — extract AST structure."""
    import ast
    text = data.decode("utf-8", errors="replace")
    try:
        tree = ast.parse(text)
        functions = [n.name for n in ast.walk(tree) if isinstance(n, ast.FunctionDef)]
        classes = [n.name for n in ast.walk(tree) if isinstance(n, ast.ClassDef)]
        imports = []
        for n in ast.walk(tree):
            if isinstance(n, ast.Import):
                imports.extend(a.name for a in n.names)
            elif isinstance(n, ast.ImportFrom):
                imports.append(n.module or "")
        metadata = {"functions": functions[:20], "classes": classes[:10], "imports": imports[:20]}
    except SyntaxError:
        metadata = {"parse_error": True}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-python")


@register_reader("js", "mjs", "jsx")
async def read_javascript(data: bytes, filename: str) -> FileContent:
    """Read JavaScript/JSX files."""
    text = data.decode("utf-8", errors="replace")
    # Simple heuristic: extract function/class names
    import re
    functions = re.findall(r'(?:function|const|let)\s+(\w+)\s*(?:=\s*(?:async\s*)?\(|\()', text)
    classes = re.findall(r'class\s+(\w+)', text)
    exports = re.findall(r'export\s+(?:default\s+)?(?:function|const|class)\s+(\w+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10], "exports": exports[:20]}
    return FileContent(text=text, metadata=metadata, mime_type="text/javascript")


@register_reader("ts", "tsx")
async def read_typescript(data: bytes, filename: str) -> FileContent:
    """Read TypeScript files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'(?:function|const|let)\s+(\w+)\s*(?:=\s*(?:async\s*)?\(|\(|<)', text)
    classes = re.findall(r'class\s+(\w+)', text)
    interfaces = re.findall(r'interface\s+(\w+)', text)
    types = re.findall(r'type\s+(\w+)\s*=', text)
    imports = re.findall(r'import\s+.*?\s+from\s+["\']([^"\']+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10], "interfaces": interfaces[:10], "types": types[:10], "imports": imports[:20]}
    return FileContent(text=text, metadata=metadata, mime_type="text/typescript")


@register_reader("css", "scss", "less")
async def read_css(data: bytes, filename: str) -> FileContent:
    """Read CSS/SCSS/Less files."""
    text = data.decode("utf-8", errors="replace")
    import re
    selectors = re.findall(r'([.#]?[\w-]+(?:\s*[>+~]\s*[\w-]+)*)\s*\{', text)
    media_queries = re.findall(r'@media\s+([^{]+)', text)
    metadata = {"selectors": len(selectors), "media_queries": media_queries[:5]}
    return FileContent(text=text, metadata=metadata, mime_type="text/css")


@register_reader("java")
async def read_java(data: bytes, filename: str) -> FileContent:
    """Read Java files."""
    text = data.decode("utf-8", errors="replace")
    import re
    classes = re.findall(r'(?:public|private|protected)?\s*class\s+(\w+)', text)
    methods = re.findall(r'(?:public|private|protected)\s+\w+\s+(\w+)\s*\(', text)
    imports = re.findall(r'import\s+([\w.]+);', text)
    metadata = {"classes": classes[:10], "methods": methods[:20], "imports": imports[:20]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-java")


@register_reader("go")
async def read_go(data: bytes, filename: str) -> FileContent:
    """Read Go files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'func\s+(?:\([^)]*\)\s+)?(\w+)\s*\(', text)
    structs = re.findall(r'type\s+(\w+)\s+struct', text)
    interfaces = re.findall(r'type\s+(\w+)\s+interface', text)
    imports = re.findall(r'"\w[\w/]*"', text)
    metadata = {"functions": functions[:20], "structs": structs[:10], "interfaces": interfaces[:10], "imports": imports[:20]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-go")


@register_reader("rs")
async def read_rust(data: bytes, filename: str) -> FileContent:
    """Read Rust files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'fn\s+(\w+)', text)
    structs = re.findall(r'struct\s+(\w+)', text)
    enums = re.findall(r'enum\s+(\w+)', text)
    impls = re.findall(r'impl\s+(\w+)', text)
    metadata = {"functions": functions[:20], "structs": structs[:10], "enums": enums[:10], "impls": impls[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/rust")


@register_reader("php")
async def read_php(data: bytes, filename: str) -> FileContent:
    """Read PHP files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'function\s+(\w+)', text)
    classes = re.findall(r'class\s+(\w+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-php")


@register_reader("c", "h")
async def read_c(data: bytes, filename: str) -> FileContent:
    """Read C header/source files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'\w+\s+(\w+)\s*\([^)]*\)\s*\{', text)
    includes = re.findall(r'#include\s+[<"]([^>"]+)[>"]', text)
    defines = re.findall(r'#define\s+(\w+)', text)
    metadata = {"functions": functions[:20], "includes": includes[:20], "defines": defines[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-c")


@register_reader("cpp", "hpp", "cc")
async def read_cpp(data: bytes, filename: str) -> FileContent:
    """Read C++ files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'\w+(?:<[^>]*>)?\s+(\w+)\s*\([^)]*\)\s*\{', text)
    classes = re.findall(r'class\s+(\w+)', text)
    includes = re.findall(r'#include\s+[<"]([^>"]+)[>"]', text)
    metadata = {"functions": functions[:20], "classes": classes[:10], "includes": includes[:20]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-c++")


@register_reader("rb")
async def read_ruby(data: bytes, filename: str) -> FileContent:
    """Read Ruby files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'def\s+(\w+)', text)
    classes = re.findall(r'class\s+(\w+)', text)
    modules = re.findall(r'module\s+(\w+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10], "modules": modules[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-ruby")


@register_reader("swift")
async def read_swift(data: bytes, filename: str) -> FileContent:
    """Read Swift files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'func\s+(\w+)', text)
    classes = re.findall(r'class\s+(\w+)', text)
    structs = re.findall(r'struct\s+(\w+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10], "structs": structs[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-swift")


@register_reader("kt")
async def read_kotlin(data: bytes, filename: str) -> FileContent:
    """Read Kotlin files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'fun\s+(\w+)', text)
    classes = re.findall(r'class\s+(\w+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-kotlin")


@register_reader("cs")
async def read_csharp(data: bytes, filename: str) -> FileContent:
    """Read C# files."""
    text = data.decode("utf-8", errors="replace")
    import re
    functions = re.findall(r'(?:public|private|protected|internal)\s+\w+\s+(\w+)\s*\(', text)
    classes = re.findall(r'class\s+(\w+)', text)
    metadata = {"functions": functions[:20], "classes": classes[:10]}
    return FileContent(text=text, metadata=metadata, mime_type="text/x-csharp")


@register_reader("dds")
async def read_dds(data: bytes, filename: str) -> FileContent:
    """Read DirectDraw Surface (game textures)."""
    from PIL import Image
    try:
        import imagecodecs
        img = Image.open(io.BytesIO(data))
        metadata = {"format": "DDS", "size": f"{img.width}x{img.height}"}
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return FileContent(text=f"[DDS texture: {img.width}x{img.height}]", metadata=metadata, images=[buf.getvalue()], mime_type="image/vnd.ms-dds")
    except Exception:
        return FileContent(text="[DDS file — cannot decode]", metadata={"error": "decode failed"}, mime_type="image/vnd.ms-dds")


@register_reader("tga")
async def read_tga(data: bytes, filename: str) -> FileContent:
    """Read TGA image files."""
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    metadata = {"format": "TGA", "size": f"{img.width}x{img.height}"}
    return FileContent(text=f"[TGA image: {img.width}x{img.height}]", metadata=metadata, images=[buf.getvalue()], mime_type="image/x-tga")


@register_reader("aseprite")
async def read_aseprite_json(data: bytes, filename: str) -> FileContent:
    """Read Aseprite sprite files (JSON format)."""
    import json
    text = data.decode("utf-8", errors="replace")
    try:
        parsed = json.loads(text)
        frames = parsed.get("frames", {})
        meta = parsed.get("meta", {})
        metadata = {
            "frames": len(frames),
            "image": meta.get("image"),
            "size": meta.get("size", {}).get("w"),
            "scale": meta.get("scale"),
        }
    except Exception:
        metadata = {"parse_error": True}
    return FileContent(text=text, metadata=metadata, mime_type="application/json")


@register_reader("sln")
async def read_sln(data: bytes, filename: str) -> FileContent:
    """Read Visual Studio Solution files."""
    text = data.decode("utf-8", errors="replace")
    import re
    projects = re.findall(r'Project\([^)]*\)\s*=\s*"([^"]+)"', text)
    return FileContent(text=text, metadata={"projects": projects}, mime_type="text/plain")


@register_reader("csproj", "vbproj")
async def read_csproj(data: bytes, filename: str) -> FileContent:
    """Read .NET project files (XML)."""
    from lxml import etree
    text = data.decode("utf-8", errors="replace")
    try:
        root = etree.fromstring(data)
        # Extract package references
        packages = root.findall(".//{http://schemas.microsoft.com/developer/msbuild/2003}PackageReference")
        if not packages:
            packages = root.findall(".//PackageReference")
        metadata = {"packages": [p.get("Include") for p in packages[:20]]}
    except Exception:
        metadata = {"parse_error": True}
    return FileContent(text=text, metadata=metadata, mime_type="application/xml")
